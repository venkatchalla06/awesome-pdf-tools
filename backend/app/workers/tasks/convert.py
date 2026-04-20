import os
import subprocess
import tempfile
import zipfile

import fitz
from PIL import Image

from app.workers.celery_app import celery_app
from app.workers.base_task import PDFBaseTask, SyncSession
from app.db.models.job import Job, JobStatus, ToolType
from app.services.storage import storage


def _is_image_pdf(pdf_path: str) -> bool:
    """Return True if the PDF has no selectable text (scanned/image-based)."""
    doc = fitz.open(pdf_path)
    has_text = any(page.get_text().strip() for page in doc)
    doc.close()
    return not has_text


def _ocr_pdf(input_path: str, output_path: str, language: str = "eng"):
    """Run OCR on a PDF, producing a searchable PDF at output_path."""
    try:
        import ocrmypdf
        ocrmypdf.ocr(
            input_path, output_path,
            language=language,
            skip_text=True,
            optimize=1,
            progress_bar=False,
            invalidate_digital_signatures=True,
        )
    except (ImportError, TypeError, AttributeError):
        import pytesseract
        doc = fitz.open(input_path)
        out_doc = fitz.open()
        for page in doc:
            mat = fitz.Matrix(300 / 72, 300 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(img, lang=language)
            new_page = out_doc.new_page(width=page.rect.width, height=page.rect.height)
            new_page.insert_image(new_page.rect, pixmap=pix)
            new_page.insert_textbox(new_page.rect, text, fontsize=1, color=(1, 1, 1), fill_opacity=0)
        doc.close()
        out_doc.save(output_path, garbage=4, deflate=True)
        out_doc.close()


def libreoffice_convert(input_path: str, output_dir: str, fmt: str = "pdf") -> str:
    """Convert using LibreOffice headless. Returns path of output file."""
    before = set(os.listdir(output_dir))

    # Each job gets its own LO user-profile so concurrent conversions don't
    # lock each other out ("source file could not be loaded" race condition).
    lo_profile = os.path.join(output_dir, ".lo_profile")
    os.makedirs(lo_profile, exist_ok=True)

    cmd = [
        "libreoffice",
        f"-env:UserInstallation=file://{lo_profile}",
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--convert-to", fmt,
        "--outdir", output_dir,
        input_path,
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=120)
    except FileNotFoundError:
        raise RuntimeError("LibreOffice is not installed. Install libreoffice-writer.")

    stderr = result.stderr.decode()
    stdout = result.stdout.decode()

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed (code {result.returncode}): {stderr[:400]}")

    # Detect the output file by directory diff (handles any name LibreOffice chose)
    after = set(os.listdir(output_dir))
    new_files = [f for f in (after - before) if f.lower().endswith(f".{fmt}")]
    if new_files:
        return os.path.join(output_dir, new_files[0])

    expected = os.path.join(output_dir, os.path.splitext(os.path.basename(input_path))[0] + f".{fmt}")
    if os.path.exists(expected):
        return expected

    raise RuntimeError(
        f"LibreOffice produced no output. stdout: {stdout[:200]} stderr: {stderr[:200]}"
    )


# ── PDF → JPG ──────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_jpg_task",
    max_retries=2, soft_time_limit=300,
)
def pdf_to_jpg_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    dpi = int(options.get("dpi", 150))

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)

        doc = fitz.open(input_path)
        total = len(doc)
        image_paths = []
        mat = fitz.Matrix(dpi / 72, dpi / 72)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_path = os.path.join(tmpdir, f"page_{i + 1:04d}.jpg")
            pix.save(img_path)
            image_paths.append(img_path)
            self.update_job(job_id, progress=5 + int((i + 1) / total * 75))

        doc.close()

        if total == 1:
            output_key = f"results/{job_id}/page_1.jpg"
            storage.upload_from_temp(image_paths[0], output_key, content_type="image/jpeg")
            output_filename = "page_1.jpg"
        else:
            zip_path = os.path.join(tmpdir, "pages.zip")
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                for p in image_paths:
                    zf.write(p, os.path.basename(p))
            output_key = f"results/{job_id}/pages.zip"
            storage.upload_from_temp(zip_path, output_key, content_type="application/zip")
            output_filename = "pages.zip"

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": output_filename}
        session.commit()

    return output_key


# ── JPG → PDF ──────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.jpg_to_pdf_task",
    max_retries=2, soft_time_limit=120,
)
def jpg_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_keys = job.input_keys
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        doc = fitz.open()

        for i, key in enumerate(input_keys):
            img_path = os.path.join(tmpdir, f"img_{i}.jpg")
            storage.download_to_temp(key, img_path)

            # Use PIL to get dimensions
            with Image.open(img_path) as img:
                w, h = img.size

            page = doc.new_page(width=w, height=h)
            page.insert_image(page.rect, filename=img_path)
            self.update_job(job_id, progress=5 + int((i + 1) / len(input_keys) * 80))

        output_path = os.path.join(tmpdir, "images.pdf")
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        output_key = f"results/{job_id}/images.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "images.pdf"}
        session.commit()

    return output_key


# ── PDF → Word ─────────────────────────────────────────────────────────────
#
# Supported Word versions via the `word_version` option:
#   "2003"        → .doc  (Word 97–2003 binary format, via LibreOffice)
#   "2007" | ""   → .docx (Office Open XML, Word 2007+)  ← default
#   "2010"–"365"  → .docx (same Open XML format, fully compatible)
#
# Strategy:
#   1. Always convert PDF → DOCX first using pdf2docx (best fidelity).
#   2. If .doc is requested, pipe the DOCX through LibreOffice headless
#      to produce a binary .doc file.

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_word_task",
    max_retries=1, soft_time_limit=360,
)
def pdf_to_word_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    # Determine output format from requested Word version
    word_version = str(options.get("word_version", "2007")).strip()
    use_doc_format = word_version == "2003"   # only 2003 needs legacy .doc
    output_ext = "doc" if use_doc_format else "docx"
    output_filename = f"converted.{output_ext}"
    content_type = (
        "application/msword"
        if use_doc_format
        else "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        docx_path  = os.path.join(tmpdir, "converted.docx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=20)

        # Step 1 — detect image-based PDF
        if _is_image_pdf(input_path):
            # OCR first, then build DOCX from extracted text (pdf2docx can't read OCR layers)
            ocr_path = os.path.join(tmpdir, "input_ocr.pdf")
            _ocr_pdf(input_path, ocr_path, language=options.get("language", "eng"))
            self.update_job(job_id, progress=40)
            import pdfplumber
            from docx import Document as DocxDocument
            from docx.shared import Pt
            worddoc = DocxDocument()
            with pdfplumber.open(ocr_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    if i > 0:
                        worddoc.add_page_break()
                    text = page.extract_text() or ""
                    for line in text.split("\n"):
                        worddoc.add_paragraph(line)
            worddoc.save(docx_path)
        else:
            # Step 1 — PDF → DOCX (pdf2docx gives best layout fidelity for digital PDFs)
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()
        self.update_job(job_id, progress=60)

        # Step 2 — DOCX → DOC via LibreOffice (Word 2003 only)
        if use_doc_format:
            doc_path = libreoffice_convert(docx_path, tmpdir, fmt="doc")
            final_path = doc_path
        else:
            final_path = docx_path

        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/{output_filename}"
        storage.upload_from_temp(final_path, output_key, content_type=content_type)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": output_filename}
        session.commit()

    return output_key


# ── Word → PDF ─────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.word_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def word_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.docx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/converted.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pdf"}
        session.commit()

    return output_key


# ── PowerPoint → PDF ───────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pptx_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def pptx_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pptx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/presentation.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "presentation.pdf"}
        session.commit()

    return output_key


# ── Excel → PDF ────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.xlsx_to_pdf_task",
    max_retries=1, soft_time_limit=300,
)
def xlsx_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.xlsx")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        output_path = libreoffice_convert(input_path, tmpdir, fmt="pdf")
        self.update_job(job_id, progress=85)

        output_key = f"results/{job_id}/spreadsheet.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "spreadsheet.pdf"}
        session.commit()

    return output_key


# ── HTML → PDF ─────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.html_to_pdf_task",
    max_retries=1, soft_time_limit=180,
)
def html_to_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.html")
        storage.download_to_temp(input_key, input_path)

        self.update_job(job_id, progress=30)

        lo_output = libreoffice_convert(input_path, tmpdir, fmt="pdf")

        output_key = f"results/{job_id}/converted.pdf"
        storage.upload_from_temp(lo_output, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pdf"}
        session.commit()

    return output_key


# ── PDF → PowerPoint ───────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_pptx_task",
    max_retries=1, soft_time_limit=300,
)
def pdf_to_pptx_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)

        # OCR image-based PDFs so slides carry selectable text
        if _is_image_pdf(input_path):
            ocr_path = os.path.join(tmpdir, "input_ocr.pdf")
            _ocr_pdf(input_path, ocr_path, language=options.get("language", "eng"))
            input_path = ocr_path
        self.update_job(job_id, progress=20)

        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        import pytesseract

        doc = fitz.open(input_path)
        total = len(doc)
        prs = Presentation()

        for i, page in enumerate(doc):
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_path = os.path.join(tmpdir, f"slide_{i:04d}.jpg")
            pix.save(img_path)

            slide_w = int(page.rect.width * 914400 / 72)
            slide_h = int(page.rect.height * 914400 / 72)
            prs.slide_width = slide_w
            prs.slide_height = slide_h

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)

            # Add invisible text box so slide text is searchable/copyable
            page_text = page.get_text().strip()
            if page_text:
                txBox = slide.shapes.add_textbox(0, 0, slide_w, slide_h)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = page_text
                run = p.runs[0]
                run.font.size = Pt(1)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            self.update_job(job_id, progress=20 + int((i + 1) / total * 70))

        doc.close()

        output_path = os.path.join(tmpdir, "converted.pptx")
        prs.save(output_path)

        output_key = f"results/{job_id}/converted.pptx"
        storage.upload_from_temp(
            output_path, output_key,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.pptx"}
        session.commit()

    return output_key


# ── PDF → Excel ────────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_xlsx_task",
    max_retries=1, soft_time_limit=300,
)
def pdf_to_xlsx_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=15)

        # OCR image-based PDFs so pdfplumber can extract text
        if _is_image_pdf(input_path):
            ocr_path = os.path.join(tmpdir, "input_ocr.pdf")
            _ocr_pdf(input_path, ocr_path, language=options.get("language", "eng"))
            input_path = ocr_path
        self.update_job(job_id, progress=30)

        import pdfplumber
        import openpyxl

        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        with pdfplumber.open(input_path) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Page {i + 1}")
                tables = page.extract_tables()
                if tables:
                    row_num = 1
                    for table in tables:
                        for row in table:
                            for col_idx, cell in enumerate(row, start=1):
                                ws.cell(row=row_num, column=col_idx, value=cell or "")
                            row_num += 1
                        row_num += 1
                else:
                    text = page.extract_text() or ""
                    for row_num, line in enumerate(text.split("\n"), start=1):
                        ws.cell(row=row_num, column=1, value=line)
                self.update_job(job_id, progress=15 + int((i + 1) / total * 70))

        output_path = os.path.join(tmpdir, "converted.xlsx")
        wb.save(output_path)

        output_key = f"results/{job_id}/converted.xlsx"
        storage.upload_from_temp(
            output_path, output_key,
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "converted.xlsx"}
        session.commit()

    return output_key


# ── Sign PDF ───────────────────────────────────────────────────────────────────
# input_keys[0] = PDF, input_keys[1] = signature image (optional)
# options: page (1-based), x%, y%, width%, height% of page dimensions

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.sign_pdf_task",
    max_retries=1, soft_time_limit=120,
)
def sign_pdf_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_keys = job.input_keys
        options = dict(job.options)

    pdf_key = input_keys[0]
    sig_key = input_keys[1] if len(input_keys) > 1 else None

    page_num = int(options.get("page", 1)) - 1
    x_pct = float(options.get("x", 60))
    y_pct = float(options.get("y", 85))
    w_pct = float(options.get("width", 25))
    h_pct = float(options.get("height", 8))

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(pdf_key, pdf_path)

        sig_path = None
        if sig_key:
            sig_path = os.path.join(tmpdir, "signature.png")
            storage.download_to_temp(sig_key, sig_path)

        self.update_job(job_id, progress=30)

        doc = fitz.open(pdf_path)
        page_idx = min(page_num, len(doc) - 1)
        page = doc[page_idx]
        pw, ph = page.rect.width, page.rect.height

        x0 = pw * x_pct / 100
        y0 = ph * y_pct / 100
        x1 = x0 + pw * w_pct / 100
        y1 = y0 + ph * h_pct / 100
        rect = fitz.Rect(x0, y0, x1, y1)

        if sig_path:
            page.insert_image(rect, filename=sig_path)
        else:
            page.draw_rect(rect, color=(0, 0, 0), width=0.5)
            page.insert_text(
                fitz.Point(x0 + 4, y0 + (y1 - y0) / 2 + 4),
                "Signed",
                fontsize=min(12, (y1 - y0) * 0.6),
                color=(0, 0.4, 0),
            )

        self.update_job(job_id, progress=70)

        output_path = os.path.join(tmpdir, "signed.pdf")
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        output_key = f"results/{job_id}/signed.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "signed.pdf"}
        session.commit()

    return output_key


# ── Fill PDF Form ──────────────────────────────────────────────────────────────
# options: {"fields": {"FieldName": "value", ...}}

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.fill_form_task",
    max_retries=1, soft_time_limit=120,
)
def fill_form_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=10)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    fields = options.get("fields", {})

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=30)

        doc = fitz.open(input_path)
        filled = 0
        for page in doc:
            for widget in page.widgets():
                if widget.field_name in fields:
                    val = fields[widget.field_name]
                    if widget.field_type == fitz.PDF_WIDGET_TYPE_CHECKBOX:
                        widget.field_value = bool(val)
                    elif widget.field_type == fitz.PDF_WIDGET_TYPE_RADIOBUTTON:
                        widget.field_value = str(val)
                    else:
                        widget.field_value = str(val)
                    widget.update()
                    filled += 1

        self.update_job(job_id, progress=70)

        output_path = os.path.join(tmpdir, "filled.pdf")
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()

        output_key = f"results/{job_id}/filled.pdf"
        storage.upload_from_temp(output_path, output_key)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "filled.pdf", "fields_filled": filled}
        session.commit()

    return output_key


# ── PDF → Markdown ─────────────────────────────────────────────────────────────

@celery_app.task(
    bind=True, base=PDFBaseTask,
    name="app.workers.tasks.convert.pdf_to_markdown_task",
    max_retries=1, soft_time_limit=300,
)
def pdf_to_markdown_task(self, job_id: str):
    self.update_job(job_id, status=JobStatus.PROCESSING, progress=5)

    with SyncSession() as session:
        job = session.get(Job, job_id)
        input_key = job.input_keys[0]
        options = dict(job.options)

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.pdf")
        storage.download_to_temp(input_key, input_path)
        self.update_job(job_id, progress=15)

        # OCR image-based PDFs so pdfplumber can extract text
        if _is_image_pdf(input_path):
            ocr_path = os.path.join(tmpdir, "input_ocr.pdf")
            _ocr_pdf(input_path, ocr_path, language=options.get("language", "eng"))
            input_path = ocr_path
        self.update_job(job_id, progress=30)

        import pdfplumber

        lines = []
        with pdfplumber.open(input_path) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                if total > 1:
                    lines.append(f"\n## Page {i + 1}\n")

                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        if not table or not table[0]:
                            continue
                        header = [str(c or "").strip() for c in table[0]]
                        lines.append("| " + " | ".join(header) + " |")
                        lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                        for row in table[1:]:
                            lines.append("| " + " | ".join(str(c or "").strip() for c in row) + " |")
                        lines.append("")

                text = page.extract_text() or ""
                for line in text.split("\n"):
                    stripped = line.strip()
                    if not stripped:
                        continue
                    if len(stripped) < 80 and stripped.isupper():
                        lines.append(f"\n### {stripped.title()}\n")
                    else:
                        lines.append(stripped)
                lines.append("")
                self.update_job(job_id, progress=15 + int((i + 1) / total * 75))

        markdown = "\n".join(lines).strip()
        output_path = os.path.join(tmpdir, "document.md")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(markdown)

        output_key = f"results/{job_id}/document.md"
        storage.upload_from_temp(output_path, output_key, content_type="text/markdown")

    with SyncSession() as session:
        job = session.get(Job, job_id)
        job.status = JobStatus.COMPLETED
        job.output_key = output_key
        job.progress = 100
        job.options = {**options, "output_filename": "document.md"}
        session.commit()

    return output_key
