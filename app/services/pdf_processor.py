import os
import tempfile
import fitz  # PyMuPDF
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
import pytesseract
from typing import List, Dict, Any
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

class PDFProcessor:
    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
    
    def merge_pdfs(self, input_files: List[str], output_path: str) -> str:
        """Merge multiple PDF files into one"""
        merger = PdfWriter()
        
        for file_path in input_files:
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                for page in reader.pages:
                    merger.add_page(page)
        
        with open(output_path, 'wb') as output_file:
            merger.write(output_file)
        
        return output_path
    
    def split_pdf(self, input_file: str, output_dir: str, mode: str = "each", every_n_pages: int = 1, custom_ranges: str = None) -> List[str]:
        """Split PDF into individual pages, every N pages, or custom ranges"""
        output_files = []
        
        with open(input_file, 'rb') as file:
            reader = PdfReader(file)
            total_pages = len(reader.pages)
            
            if mode == "each":
                for page_num in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[page_num])
                    output_path = os.path.join(output_dir, f"page_{page_num + 1}.pdf")
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    output_files.append(output_path)
            
            elif mode == "every":
                if every_n_pages < 1:
                    every_n_pages = 1
                for i in range(0, total_pages, every_n_pages):
                    writer = PdfWriter()
                    end_page = min(i + every_n_pages, total_pages)
                    for page_num in range(i, end_page):
                        writer.add_page(reader.pages[page_num])
                    output_path = os.path.join(output_dir, f"pages_{i + 1}_to_{end_page}.pdf")
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    output_files.append(output_path)
            
            elif mode == "custom":
                if custom_ranges:
                    ranges = [r.strip() for r in custom_ranges.split(',') if r.strip()]
                    for r in ranges:
                        writer = PdfWriter()
                        if '-' in r:
                            parts = r.split('-')
                            try:
                                start = max(1, int(parts[0].strip()))
                                end = min(total_pages, int(parts[1].strip()))
                                if start <= end:
                                    for page_num in range(start - 1, end):
                                        writer.add_page(reader.pages[page_num])
                                    output_path = os.path.join(output_dir, f"range_{start}_to_{end}.pdf")
                            except ValueError:
                                continue
                        else:
                            try:
                                page_num = int(r.strip())
                                if 1 <= page_num <= total_pages:
                                    writer.add_page(reader.pages[page_num - 1])
                                    output_path = os.path.join(output_dir, f"page_{page_num}.pdf")
                            except ValueError:
                                continue
                        
                        if len(writer.pages) > 0:
                            base_path = output_path
                            counter = 1
                            while os.path.exists(output_path):
                                output_path = base_path.replace('.pdf', f'_{counter}.pdf')
                                counter += 1
                                
                            with open(output_path, 'wb') as output_file:
                                writer.write(output_file)
                            output_files.append(output_path)
            
            if not output_files:
                for page_num in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[page_num])
                    output_path = os.path.join(output_dir, f"page_{page_num + 1}.pdf")
                    with open(output_path, 'wb') as output_file:
                        writer.write(output_file)
                    output_files.append(output_path)
        
        return output_files
    
    def compress_pdf(self, input_file: str, output_path: str, quality: str = "medium") -> str:
        """Compress PDF using Ghostscript"""
        quality_settings = {
            "low": "/ebook",
            "medium": "/printer", 
            "high": "/prepress"
        }
        
        gs_quality = quality_settings.get(quality, "/printer")
        
        args = [
            "gs",
            "-dNOPAUSE", "-dBATCH", "-dSAFER",
            "-sDEVICE=pdfwrite",
            f"-dPDFSETTINGS={gs_quality}",
            "-dCompatibilityLevel=1.4",
            "-dEmbedAllFonts=true",
            "-dSubsetFonts=true",
            "-dColorImageDownsampleType=/Bicubic",
            "-dColorImageResolution=150",
            "-dGrayImageDownsampleType=/Bicubic",
            "-dGrayImageResolution=150",
            "-dMonoImageDownsampleType=/Bicubic",
            "-dMonoImageResolution=150",
            f"-sOutputFile={output_path}",
            input_file
        ]
        
        import ghostscript  # requires native libgs — install via: brew install ghostscript
        ghostscript.Ghostscript(*args)
        return output_path
    
    def pdf_to_images(self, input_file: str, output_dir: str, dpi: int = 300) -> List[str]:
        """Convert PDF pages to images"""
        doc = fitz.open(input_file)
        image_files = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72))
            output_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
            pix.save(output_path)
            image_files.append(output_path)
        
        doc.close()
        return image_files
    
    def images_to_pdf(self, image_files: List[str], output_path: str) -> str:
        """Convert images to PDF"""
        images = [Image.open(img) for img in image_files]
        
        if images:
            images[0].save(
                output_path, "PDF", resolution=100.0,
                save_all=True, append_images=images[1:]
            )
        
        return output_path
    
    def ocr_pdf(self, input_file: str, output_path: str) -> str:
        """Perform OCR on scanned PDF"""
        # Convert PDF to images
        temp_dir = tempfile.mkdtemp()
        image_files = self.pdf_to_images(input_file, temp_dir)
        
        # Create searchable PDF
        doc = fitz.open()
        
        for img_path in image_files:
            # Perform OCR
            text = pytesseract.image_to_string(Image.open(img_path))
            
            # Create PDF page with OCR text
            page = doc.new_page()
            
            # Add OCR text as invisible layer
            if text.strip():
                rect = page.rect
                page.insert_text((50, 50), text, fontsize=1)  # Tiny invisible text
            
            # Add original image
            page.insert_image(rect, filename=img_path)
        
        doc.save(output_path)
        doc.close()
        
        # Cleanup
        for img in image_files:
            os.unlink(img)
        os.rmdir(temp_dir)
        
        return output_path
    
    def rotate_pdf(self, input_file: str, output_path: str, rotation: int) -> str:
        """Rotate PDF pages"""
        doc = fitz.open(input_file)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page.set_rotation(rotation)
        
        doc.save(output_path)
        doc.close()
        return output_path
    
    def add_watermark(self, input_file: str, output_path: str, watermark_text: str) -> str:
        """Add text watermark to PDF"""
        doc = fitz.open(input_file)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Add watermark text
            rect = page.rect
            text_width = len(watermark_text) * 10
            
            page.insert_text(
                ((rect.width - text_width) / 2, rect.height / 2),
                watermark_text,
                fontsize=40,
                color=(0.8, 0.8, 0.8),  # Light gray
                rotate=45
            )
        
        doc.save(output_path)
        doc.close()
        return output_path
    
    def pdf_to_powerpoint(self, input_file: str, output_path: str, dpi: int = 150) -> str:
        """Convert PDF pages to PowerPoint slides (one slide per page)"""
        doc = fitz.open(input_file)
        prs = Presentation()

        # Use widescreen 16:9 slide dimensions
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # completely blank layout

        with tempfile.TemporaryDirectory() as tmp_dir:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
                img_path = os.path.join(tmp_dir, f"page_{page_num}.png")
                pix.save(img_path)

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    img_path,
                    left=Inches(0),
                    top=Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

        doc.close()
        prs.save(output_path)
        return output_path

    def add_page_numbers(self, input_file: str, output_path: str) -> str:
        """Add page numbers to PDF"""
        doc = fitz.open(input_file)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            rect = page.rect
            
            page.insert_text(
                (rect.width - 50, rect.height - 20),
                f"Page {page_num + 1}",
                fontsize=12,
                color=(0, 0, 0)
            )
        
        doc.save(output_path)
        doc.close()
        return output_path